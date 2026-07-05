#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère le PowerPoint :
    « La cybersécurité : protéger ses données à l'ère du numérique »

Public non spécialiste — durée cible 2 h.
Sortie : cybersecurite_presentation.pptx (format 16:9).

Dépendances : python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# --------------------------------------------------------------------------
# Palette de couleurs
# --------------------------------------------------------------------------
NAVY   = RGBColor(0x0F, 0x2A, 0x47)   # bleu nuit (fond des sections)
NAVY2  = RGBColor(0x14, 0x35, 0x59)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # bleu accent
TEAL   = RGBColor(0x1F, 0x9E, 0x8F)   # vert/teal = protection
RED    = RGBColor(0xC0, 0x39, 0x2B)   # rouge = menace / danger
ORANGE = RGBColor(0xE0, 0x6C, 0x2B)   # orange = vigilance
LIGHT  = RGBColor(0xF2, 0xF6, 0xFB)   # fond clair
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD6, 0xDF, 0xEA)
DARK   = RGBColor(0x1B, 0x2A, 0x3A)   # texte principal
GREY   = RGBColor(0x5A, 0x6B, 0x7B)   # texte secondaire
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

_page = 0  # compteur de diapositives pour le pied de page


# --------------------------------------------------------------------------
# Fonctions utilitaires
# --------------------------------------------------------------------------
def new_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, l, t, w, h, color, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def rounded(slide, l, t, w, h, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         first=False, space_after=6, space_before=0, font=FONT, line_spacing=1.05):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def footer(slide, section=""):
    """Pied de page : num. de diapo + rappel du titre."""
    global _page
    _page += 1
    tb, tf = textbox(slide, Inches(0.5), Inches(7.02), Inches(9), Inches(0.4))
    para(tf, "Cybersécurité — protéger ses données à l'ère du numérique",
         9, GREY, first=True)
    tb2, tf2 = textbox(slide, Inches(11.5), Inches(7.02), Inches(1.3), Inches(0.4))
    para(tf2, str(_page), 10, GREY, align=PP_ALIGN.RIGHT, first=True, bold=True)


def content_header(slide, title, kicker=None, accent=BLUE):
    """Barre supérieure + titre pour les diapos de contenu."""
    rect(slide, 0, 0, SW, SH, LIGHT)                 # fond clair
    rect(slide, 0, 0, SW, Inches(0.18), accent)      # fine barre couleur
    top = Inches(0.5)
    if kicker:
        tb, tf = textbox(slide, Inches(0.6), Inches(0.42), Inches(11), Inches(0.4))
        para(tf, kicker.upper(), 12, accent, bold=True, first=True)
        top = Inches(0.85)
    tb, tf = textbox(slide, Inches(0.6), top, Inches(12.1), Inches(0.9))
    para(tf, title, 30, NAVY, bold=True, first=True, font=FONT_H)
    # trait sous le titre
    rect(slide, Inches(0.62), top + Inches(0.92), Inches(1.4), Pt(3), accent)


# --------------------------------------------------------------------------
# Modèles de diapositives
# --------------------------------------------------------------------------
def slide_title(title, subtitle, meta):
    s = new_slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, SW, SH, NAVY)
    # bandes décoratives
    rect(s, 0, Inches(5.55), SW, Inches(0.09), BLUE)
    rect(s, 0, Inches(5.7), SW, Inches(0.05), TEAL)
    # petit "cadenas" stylisé (rectangle + anse)
    rounded(s, Inches(0.9), Inches(0.85), Inches(0.7), Inches(0.55), TEAL)
    an = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(1.02), Inches(0.55),
                            Inches(0.46), Inches(0.5), )
    an.fill.background()
    an.line.color.rgb = TEAL
    an.line.width = Pt(3)
    an.shadow.inherit = False
    tb, tf = textbox(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.6))
    para(tf, "LA CYBERSÉCURITÉ", 15, TEAL, bold=True, first=True, space_after=6)
    para(tf, title, 54, WHITE, bold=True, font=FONT_H, space_after=6, line_spacing=1.0)
    para(tf, subtitle, 24, RGBColor(0xC7, 0xD6, 0xE6), space_before=4)
    tb2, tf2 = textbox(s, Inches(0.6), Inches(5.95), Inches(12), Inches(1.0))
    para(tf2, meta, 15, RGBColor(0x9F, 0xB3, 0xC9), first=True)
    return s


def slide_section(number, title, subtitle=""):
    s = new_slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, Inches(0.28), SH, TEAL)
    tb, tf = textbox(s, Inches(0.95), Inches(1.7), Inches(11.5), Inches(1.2))
    para(tf, number, 90, NAVY2 if False else RGBColor(0x35, 0x59, 0x7E),
         bold=True, first=True, font=FONT_H)
    # numéro plus lisible en accent
    tb.text_frame.paragraphs[0].runs[0].font.color.rgb = TEAL
    tb2, tf2 = textbox(s, Inches(1.0), Inches(3.2), Inches(11.4), Inches(2.2))
    para(tf2, title, 42, WHITE, bold=True, first=True, font=FONT_H, line_spacing=1.0)
    if subtitle:
        para(tf2, subtitle, 20, RGBColor(0xB9, 0xC8, 0xDA), space_before=10)
    rect(s, Inches(1.03), Inches(3.05), Inches(1.6), Pt(4), TEAL)
    return s


def bullets(tf, items, size=20, color=DARK, accent=BLUE, first_para=True,
            space_after=11):
    """items : liste de str ou de tuples (texte, niveau)."""
    started = not first_para
    for it in items:
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if (first_para and not started) else tf.add_paragraph()
        started = True
        p.space_after = Pt(space_after if lvl == 0 else 4)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        indent = "" if lvl == 0 else "      "
        mk = p.add_run()
        mk.text = indent + ("●  " if lvl == 0 else "–  ")
        mk.font.size = Pt(size if lvl == 0 else size - 4)
        mk.font.color.rgb = accent if lvl == 0 else GREY
        mk.font.bold = True
        mk.font.name = FONT
        rn = p.add_run()
        rn.text = text
        rn.font.size = Pt(size if lvl == 0 else size - 4)
        rn.font.color.rgb = color if lvl == 0 else GREY
        rn.font.name = FONT
        rn.font.bold = False


def slide_bullets(title, items, kicker=None, accent=BLUE, size=20, note=None):
    s = new_slide()
    content_header(s, title, kicker, accent)
    tb, tf = textbox(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.7))
    bullets(tf, items, size=size, accent=accent)
    if note:
        rounded(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.7),
                RGBColor(0xEA, 0xF1, 0xF9))
        tb2, tf2 = textbox(s, Inches(0.95), Inches(6.1), Inches(11.5), Inches(0.6),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tf2, note, 15, NAVY, italic=True, first=True)
    footer(s)
    return s


def slide_two_col(title, left_head, left_items, right_head, right_items,
                  kicker=None, accent=BLUE, left_color=RED, right_color=TEAL):
    s = new_slide()
    content_header(s, title, kicker, accent)
    top = Inches(2.0)
    colw = Inches(5.85)
    gap = Inches(0.4)
    # colonne gauche
    rounded(s, Inches(0.6), top, colw, Inches(4.6), CARD, BORDER)
    rect(s, Inches(0.6), top, colw, Inches(0.62), left_color)
    tb, tf = textbox(s, Inches(0.8), top + Inches(0.06), colw - Inches(0.4),
                     Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, left_head, 18, WHITE, bold=True, first=True)
    tb, tf = textbox(s, Inches(0.85), top + Inches(0.8), colw - Inches(0.5), Inches(3.6))
    bullets(tf, left_items, size=16, accent=left_color, space_after=8)
    # colonne droite
    lx = Inches(0.6) + colw + gap
    rounded(s, lx, top, colw, Inches(4.6), CARD, BORDER)
    rect(s, lx, top, colw, Inches(0.62), right_color)
    tb, tf = textbox(s, lx + Inches(0.2), top + Inches(0.06), colw - Inches(0.4),
                     Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, right_head, 18, WHITE, bold=True, first=True)
    tb, tf = textbox(s, lx + Inches(0.25), top + Inches(0.8), colw - Inches(0.5), Inches(3.6))
    bullets(tf, right_items, size=16, accent=right_color, space_after=8)
    footer(s)
    return s


def slide_stat(big, caption, sub="", accent=BLUE):
    s = new_slide()
    rect(s, 0, 0, SW, SH, accent)
    rect(s, 0, 0, SW, SH, accent)
    tb, tf = textbox(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.2),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, 110, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         font=FONT_H, line_spacing=0.9)
    tb2, tf2 = textbox(s, Inches(1.2), Inches(4.5), Inches(10.9), Inches(1.6),
                       anchor=MSO_ANCHOR.TOP)
    para(tf2, caption, 26, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         line_spacing=1.05)
    if sub:
        para(tf2, sub, 15, RGBColor(0xE7, 0xEF, 0xF7), align=PP_ALIGN.CENTER,
             italic=True, space_before=8)
    return s


def slide_cards(title, cards, kicker=None, accent=BLUE, cols=3):
    """cards : liste de (titre, description) ou (titre, description, couleur)."""
    s = new_slide()
    content_header(s, title, kicker, accent)
    margin = Inches(0.6)
    gutter = Inches(0.3)
    total_w = SW - 2 * margin
    card_w = int((total_w - (cols - 1) * gutter) / cols)
    rows = math.ceil(len(cards) / cols)
    top0 = Inches(2.05)
    avail_h = Inches(4.7)
    card_h = int((avail_h - (rows - 1) * gutter) / rows)
    for i, c in enumerate(cards):
        if len(c) == 3:
            ctitle, desc, col = c
        else:
            ctitle, desc = c
            col = accent
        r = i // cols
        cc = i % cols
        l = margin + cc * (card_w + gutter)
        t = top0 + r * (card_h + gutter)
        rounded(s, l, t, card_w, card_h, CARD, BORDER)
        rect(s, l, t, Inches(0.12), card_h, col)          # barre latérale couleur
        tb, tf = textbox(s, l + Inches(0.3), t + Inches(0.14),
                         card_w - Inches(0.45), card_h - Inches(0.25))
        para(tf, ctitle, 16, NAVY, bold=True, first=True, space_after=4)
        if desc:
            para(tf, desc, 12.5, GREY, line_spacing=1.0)
    footer(s)
    return s


def slide_checklist(title, items, kicker=None, accent=TEAL, cols=2, note=None):
    s = new_slide()
    content_header(s, title, kicker, accent)
    margin = Inches(0.7)
    colw = Inches(5.9)
    gap = Inches(0.4)
    per_col = math.ceil(len(items) / cols)
    top0 = 2.05
    line_h = 0.62
    for i, it in enumerate(items):
        col = i // per_col
        row = i % per_col
        l = margin + col * (colw + gap)
        t = Inches(top0 + row * line_h)
        # pastille numérotée
        badge = rounded(s, l, t, Inches(0.42), Inches(0.42), accent)
        tbb, tfb = textbox(s, l, t - Inches(0.02), Inches(0.42), Inches(0.46),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tfb, str(i + 1), 14, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
        tb, tf = textbox(s, l + Inches(0.55), t - Inches(0.04), colw - Inches(0.55),
                         Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, it, 16, DARK, first=True, line_spacing=1.0)
    if note:
        rounded(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.62), NAVY)
        tb2, tf2 = textbox(s, Inches(0.95), Inches(6.36), Inches(11.5), Inches(0.6),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tf2, note, 15, WHITE, bold=True, first=True)
    footer(s)
    return s


def slide_quiz(title, qa, kicker="ON TESTE ENSEMBLE", accent=ORANGE):
    """qa : liste de (question, réponse)."""
    s = new_slide()
    content_header(s, title, kicker, accent)
    top = 2.05
    for i, (q, a) in enumerate(qa):
        t = Inches(top + i * 1.15)
        rounded(s, Inches(0.7), t, Inches(11.9), Inches(1.0), CARD, BORDER)
        rect(s, Inches(0.7), t, Inches(0.12), Inches(1.0), accent)
        tb, tf = textbox(s, Inches(1.0), t + Inches(0.08), Inches(11.4), Inches(0.9),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "❔ " + q, 16, NAVY, bold=True, first=True, space_after=3)
        para(tf, "➡ " + a, 14, TEAL, italic=True)
    footer(s)
    return s


# ==========================================================================
#  CONSTRUCTION DE LA PRÉSENTATION
# ==========================================================================

# --- Ouverture -----------------------------------------------------------
slide_title(
    "Protéger ses données\nà l'ère du numérique",
    "Comprendre les risques et adopter les bons réflexes",
    "Atelier de sensibilisation  ·  Durée : 2 heures  ·  Public : tout le monde, aucun prérequis technique",
)

slide_bullets(
    "Les objectifs de cet atelier",
    [
        "Comprendre pourquoi la cybersécurité nous concerne TOUS",
        "Reconnaître les principales menaces du quotidien",
        "Adopter des réflexes simples, concrets et gratuits",
        "Savoir réagir si l'on est victime d'une attaque",
        "Repartir avec une check-list et des ressources fiables",
    ],
    kicker="Bienvenue",
    accent=TEAL,
    note="Aucune compétence technique requise : on parle de la vraie vie, pas de code.",
)

slide_cards(
    "Au programme des 2 heures",
    [
        ("Partie 1 — Comprendre", "Les bases + le panorama des menaces", BLUE),
        ("Pause", "~10 minutes au milieu de la séance", GREY),
        ("Partie 2 — Agir", "Se protéger, cas concrets, réagir", TEAL),
        ("Interactif", "Quiz, questions et échanges tout du long", ORANGE),
        ("Check-list finale", "10 habitudes à emporter chez soi", TEAL),
        ("Questions", "On prend le temps d'y répondre", BLUE),
    ],
    kicker="Déroulé",
    accent=BLUE,
    cols=3,
)

slide_bullets(
    "Un petit sondage pour démarrer",
    [
        "Qui a déjà reçu un SMS ou un e-mail qui semblait suspect ?",
        "Qui utilise (parfois) le même mot de passe sur plusieurs sites ?",
        "Qui a déjà hésité avant de cliquer sur un lien ?",
        "Qui a déjà eu un proche victime d'une arnaque en ligne ?",
    ],
    kicker="À main levée",
    accent=ORANGE,
    note="Si vous avez levé la main au moins une fois : vous êtes exactement au bon endroit.",
)

# --- PARTIE 1 : COMPRENDRE ----------------------------------------------
slide_section("01", "Comprendre le monde numérique", "De quoi parle-t-on, et pourquoi c'est important ?")

slide_bullets(
    "Notre vie est devenue numérique",
    [
        "Banque, impôts, santé, travail, école, loisirs, famille…",
        "Le smartphone est devenu le prolongement de nous-mêmes",
        "Nous produisons des données en permanence, souvent sans le savoir",
        "Le confort du « tout en ligne » a un revers : de nouveaux risques",
    ],
    kicker="Le contexte",
    accent=BLUE,
    note="Plus notre vie est connectée, plus la surface exposée est grande.",
)

slide_stat(
    "39 s",
    "En moyenne, une attaque informatique aurait lieu toutes les 39 secondes dans le monde",
    "Chiffre illustratif souvent cité (étude Univ. du Maryland) — l'ordre de grandeur compte plus que la précision.",
    accent=RED,
)

slide_stat(
    "≈ 90 %",
    "des cyberattaques commencent par un simple e-mail (hameçonnage)",
    "La porte d'entrée n°1, c'est nous — pas une faille technique obscure.",
    accent=BLUE,
)

slide_bullets(
    "Qu'est-ce que la cybersécurité ?",
    [
        "L'ensemble des moyens pour protéger nos appareils, nos comptes et nos données",
        "Objectif : empêcher l'accès, le vol, la modification ou la destruction non autorisés",
        ("La bonne image : la sécurité d'une maison", 0),
        ("porte, serrure, alarme = les outils techniques", 1),
        ("mais surtout : fermer à clé, ne pas ouvrir à un inconnu = les comportements", 1),
    ],
    kicker="Définition",
    accent=BLUE,
    note="80 % de la sécurité, ce sont des habitudes. La technique ne fait pas tout.",
)

slide_cards(
    "Une « donnée », c'est quoi au juste ?",
    [
        ("Identité", "Nom, date de naissance, n° de sécu", BLUE),
        ("Coordonnées", "Adresse, téléphone, e-mail", BLUE),
        ("Accès", "Identifiants et mots de passe", RED),
        ("Argent", "Coordonnées bancaires, cartes", RED),
        ("Vie privée", "Photos, messages, contacts", TEAL),
        ("Traces", "Localisation, historique, achats", ORANGE),
    ],
    kicker="Vos informations",
    accent=BLUE,
    cols=3,
)

slide_bullets(
    "Pourquoi VOS données ont de la valeur",
    [
        "« Je n'ai rien à cacher » ≠ « je n'ai rien à protéger »",
        "Les données volées se revendent (marché noir, « dark web »)",
        "Elles servent à l'usurpation d'identité et à la fraude bancaire",
        "Elles alimentent des arnaques ciblées, très crédibles",
        "Rappel : quand un service est gratuit, la donnée est souvent le produit",
    ],
    kicker="La cible, c'est vous",
    accent=RED,
    note="On ne vole pas que les célébrités : les particuliers sont des cibles faciles et rentables.",
)

slide_cards(
    "Les 3 piliers de la sécurité",
    [
        ("Confidentialité", "Seules les bonnes personnes accèdent à l'information", BLUE),
        ("Intégrité", "L'information n'est pas modifiée à votre insu", TEAL),
        ("Disponibilité", "Vous y avez accès quand vous en avez besoin", ORANGE),
    ],
    kicker="Le B.A.-BA",
    accent=BLUE,
    cols=3,
)

slide_cards(
    "Qui sont les attaquants ?",
    [
        ("Cybercriminels", "Motivés par l'argent (l'immense majorité)", RED),
        ("Arnaqueurs", "Opportunistes, en masse et automatisés", ORANGE),
        ("Hacktivistes", "Militent pour une cause", BLUE),
        ("Espionnage / États", "Cibles stratégiques, entreprises", NAVY),
        ("Proches malveillants", "Ex-conjoint, entourage, curieux", TEAL),
        ("L'erreur humaine", "Pas malveillant, mais fréquent", GREY),
    ],
    kicker="Le portrait-robot",
    accent=RED,
    cols=3,
)
slide_bullets(
    "Le mythe du génie encapuchonné",
    [
        "L'image du hacker seul dans le noir est trompeuse",
        "La réalité : une industrie organisée, avec des « services » et de la sous-traitance",
        "Des attaques automatisées qui ratissent large, 24 h/24",
        "On ne vous vise pas forcément VOUS : on vise tout le monde, et on attend qui tombe",
    ],
    kicker="Réalité vs fiction",
    accent=NAVY,
    note="Bonne nouvelle : contre des attaques automatisées, des réflexes simples suffisent souvent.",
)

# --- PARTIE 2 : LES MENACES ---------------------------------------------
slide_section("02", "Les menaces les plus courantes", "Reconnaître l'ennemi pour mieux s'en protéger")

slide_cards(
    "Panorama des menaces",
    [
        ("Logiciels malveillants", "Virus, rançongiciels, espions", RED),
        ("Hameçonnage", "Faux messages qui piègent", ORANGE),
        ("Ingénierie sociale", "Manipulation de l'humain", ORANGE),
        ("Vol de mots de passe", "Comptes piratés en cascade", RED),
        ("Fuites de données", "Vos infos exposées", BLUE),
        ("Arnaques en ligne", "Faux sites, faux support", ORANGE),
    ],
    kicker="Vue d'ensemble",
    accent=RED,
    cols=3,
)

slide_bullets(
    "Les logiciels malveillants (« malwares »)",
    [
        "Des programmes conçus pour nuire ou espionner",
        ("Les grandes familles :", 0),
        ("Virus / vers : se propagent d'un appareil à l'autre", 1),
        ("Cheval de Troie : se cache dans un logiciel d'apparence normale", 1),
        ("Espiogiciel (spyware) : surveille ce que vous faites", 1),
        ("Comment ils arrivent : pièce jointe, téléchargement, fausse mise à jour, clé USB",  0),
    ],
    kicker="Menace n°1",
    accent=RED,
    size=19,
)

slide_bullets(
    "Focus : les rançongiciels (ransomware)",
    [
        "Ils chiffrent (verrouillent) vos fichiers, puis exigent une rançon",
        "Photos, documents, souvenirs : tout devient inaccessible d'un coup",
        "Ont paralysé des hôpitaux, des mairies, des écoles, des entreprises",
        "Payer ne garantit RIEN — et encourage les criminels",
        "La meilleure parade existe et elle est simple : la sauvegarde",
    ],
    kicker="La hantise du moment",
    accent=RED,
    note="On en reparle en Partie 2 : une bonne sauvegarde réduit un drame à un simple contretemps.",
)

slide_bullets(
    "Le hameçonnage (phishing) — le plus fréquent",
    [
        "Un faux message qui imite un organisme de confiance",
        ("Banque, impôts, CAF, Ameli, La Poste, Netflix, un livreur…", 1),
        "But : vous faire cliquer, puis saisir vos identifiants ou payer",
        ("Il se décline sur tous les canaux :", 0),
        ("par e-mail (phishing), par SMS (« smishing »), par téléphone (« vishing »)", 1),
    ],
    kicker="À connaître par cœur",
    accent=ORANGE,
    size=19,
)

slide_cards(
    "Anatomie d'un message piégé",
    [
        ("1. Expéditeur douteux", "Adresse bizarre, imitation approximative", RED),
        ("2. Urgence / peur", "« Votre compte va être suspendu ! »", RED),
        ("3. Formulation étrange", "Fautes, ton inhabituel, « Cher client »", ORANGE),
        ("4. Lien trompeur", "Le texte affiché ≠ la vraie adresse", ORANGE),
        ("5. Demande sensible", "Mot de passe, code, coordonnées bancaires", RED),
        ("6. Pièce jointe", "Facture ou colis « à ouvrir vite »", ORANGE),
    ],
    kicker="Les 6 signaux d'alerte",
    accent=ORANGE,
    cols=3,
)

slide_bullets(
    "L'ingénierie sociale : manipuler l'humain",
    [
        "La technique la plus efficace ne vise pas la machine… mais vous",
        ("Les leviers psychologiques exploités :", 0),
        ("l'urgence, la peur, l'autorité, l'appât du gain, la curiosité, la gentillesse", 1),
        "Exemples : faux service technique, « faux patron », faux conseiller bancaire",
        "Un attaquant préfère vous demander la clé plutôt que de forcer la porte",
    ],
    kicker="La faille, c'est l'humain",
    accent=ORANGE,
    note="Le doute est votre meilleur allié : un organisme sérieux ne vous mettra jamais la pression.",
)

slide_cards(
    "Les arnaques que vous croiserez",
    [
        ("Colis en attente", "SMS « payez 1,99 € de frais »", ORANGE),
        ("Faux support", "« Votre PC est infecté, appelez… »", RED),
        ("Conseiller bancaire", "« Un virement suspect, confirmez »", RED),
        ("Bonnes affaires", "Site qui imite une grande marque", ORANGE),
        ("Arnaque sentimentale", "Belle rencontre qui finit par un virement", ORANGE),
        ("Faux remboursement", "« Les impôts vous doivent 200 € »", RED),
    ],
    kicker="Ça n'arrive pas qu'aux autres",
    accent=RED,
    cols=3,
)

slide_two_col(
    "Le nerf de la guerre : les mots de passe",
    "❌ Ce que font les attaquants",
    [
        "Réutilisent les fuites : un mot de passe volé, testé partout",
        "Devinent : prénom, date de naissance, « 123456 »",
        "Force brute : des millions d'essais automatiques",
        "Rachètent des listes d'identifiants tout prêts",
    ],
    "🎯 Pourquoi ça marche",
    [
        "On réutilise le même mot de passe partout",
        "On choisit des mots de passe trop courts / trop simples",
        "Une seule fuite ouvre alors TOUS nos comptes",
        "Effet domino : e-mail piraté = tout le reste tombe",
    ],
    kicker="Vol d'identifiants",
    accent=RED,
)

slide_bullets(
    "Les fuites de données (« data breaches »)",
    [
        "Même les grandes entreprises se font pirater leurs bases clients",
        "Vos e-mail, mots de passe, numéros peuvent circuler à votre insu",
        "Vous n'y pouvez rien directement… mais vous pouvez limiter les dégâts",
        "Outil gratuit : « Have I Been Pwned » pour vérifier si un e-mail a fuité",
    ],
    kicker="Hors de votre contrôle",
    accent=BLUE,
    note="D'où l'importance d'un mot de passe UNIQUE par site : une fuite reste alors isolée.",
)

slide_bullets(
    "Le Wi-Fi public : pratique mais risqué",
    [
        "Gare, café, hôtel, aéroport : réseaux ouverts et partagés",
        "Un tiers mal intentionné peut parfois « écouter » le trafic",
        "Faux réseaux : « Free_WiFi_Gratuit » créé par un pirate",
        ("Les bons réflexes :", 0),
        ("éviter banque et achats, préférer le partage de connexion (4G/5G), utiliser un VPN", 1),
    ],
    kicker="En déplacement",
    accent=ORANGE,
    size=19,
)

slide_quiz(
    "Quiz express — vrai ou faux ?",
    [
        ("Le petit cadenas « https » garantit que le site est honnête.",
         "FAUX : il garantit que la connexion est chiffrée, pas que le site est fiable."),
        ("Un mot de passe très compliqué est incassable.",
         "FAUX : s'il fuite dans une base de données, sa complexité n'y change rien."),
        ("Mon smartphone n'a besoin d'aucune précaution.",
         "FAUX : c'est l'appareil le plus personnel… donc le plus sensible."),
    ],
)

# --- PAUSE ---------------------------------------------------------------
s = slide_section("⏸", "Pause", "On se retrouve dans 10 minutes — n'hésitez pas à noter vos questions !")

# --- PARTIE 3 : SE PROTÉGER ---------------------------------------------
slide_section("03", "Se protéger au quotidien", "Des réflexes simples, gratuits et efficaces")

slide_cards(
    "Si vous ne retenez que 3 choses",
    [
        ("Mots de passe solides", "Longs, uniques, + un gestionnaire", TEAL),
        ("Mises à jour", "Toujours faites, de préférence auto", TEAL),
        ("Vigilance", "Réfléchir 3 secondes avant de cliquer", TEAL),
    ],
    kicker="L'essentiel",
    accent=TEAL,
    cols=3,
)

slide_two_col(
    "Les mots de passe : erreurs vs bonnes pratiques",
    "❌ À éviter absolument",
    [
        "123456, azerty, motdepasse, 0000",
        "Prénom + date de naissance",
        "Le même mot de passe partout",
        "Noté sur un post-it collé à l'écran",
        "Partagé par SMS ou par e-mail",
    ],
    "✅ À adopter",
    [
        "Au moins 12 à 14 caractères",
        "Une « phrase de passe » facile à retenir",
        "Un mot de passe DIFFÉRENT par site important",
        "Stocké dans un gestionnaire dédié",
        "La double authentification en plus",
    ],
    kicker="Le socle de tout",
    accent=TEAL,
)

slide_bullets(
    "Fabriquer un bon mot de passe",
    [
        "La méthode gagnante : la « phrase de passe »",
        ("Prenez une phrase absurde et facile à retenir…", 1),
        ("« MonChatGrisMange3Croquettes! » → long, unique, mémorisable", 1),
        "Plus c'est long, mieux c'est (la longueur prime sur la complexité)",
        "Une variante par site, ou mieux : laissez le gestionnaire les générer",
    ],
    kicker="En pratique",
    accent=TEAL,
    note="Un mot de passe de 4 mots est plus solide ET plus simple à retenir que « P@ssw0rd! ».",
)

slide_bullets(
    "Le gestionnaire de mots de passe",
    [
        "Un coffre-fort numérique qui retient tous vos mots de passe",
        "Vous ne mémorisez plus qu'UN seul mot de passe maître (très solide)",
        "Il en génère de longs, uniques et différents pour chaque site",
        "Il les remplit automatiquement : plus pratique ET plus sûr",
        "Exemples grand public : Bitwarden, KeePass, ou celui de votre navigateur",
    ],
    kicker="L'outil qui change tout",
    accent=TEAL,
    note="C'est LE meilleur retour sur investissement en cybersécurité pour un particulier.",
)

slide_bullets(
    "La double authentification (2FA / A2F)",
    [
        "Un 2ᵉ facteur en plus du mot de passe : un code, une appli, une clé",
        "Même si l'on vole votre mot de passe, le compte reste verrouillé",
        ("Les formes courantes :", 0),
        ("code par SMS, application (Google/Microsoft Authenticator), clé physique", 1),
        "À activer EN PRIORITÉ sur : e-mail, banque, réseaux sociaux",
    ],
    kicker="Votre meilleure protection",
    accent=TEAL,
    note="C'est comme un 2ᵉ verrou : un peu moins pratique, infiniment plus sûr.",
)

slide_bullets(
    "Les mises à jour : ne pas remettre à demain",
    [
        "Une mise à jour corrige souvent des failles de sécurité connues",
        "Reporter une mise à jour, c'est laisser une porte ouverte",
        "Activez les mises à jour AUTOMATIQUES quand c'est possible",
        "Cela concerne tout : téléphone, ordinateur, applis, box, TV, objets connectés",
    ],
    kicker="Le geste gratuit",
    accent=TEAL,
    note="« Installer maintenant » plutôt que « Me le rappeler demain » : votre futur vous remerciera.",
)

slide_bullets(
    "Antivirus, pare-feu : que faut-il vraiment ?",
    [
        "Sur ordinateur (Windows) : un antivirus reste utile",
        "Windows Defender, intégré et gratuit, suffit pour la plupart des gens",
        "Le pare-feu (activé par défaut) filtre les connexions : laissez-le allumé",
        "Aucun logiciel ne remplace la vigilance : l'antivirus est un filet, pas un bouclier",
    ],
    kicker="Faire simple",
    accent=BLUE,
    note="Méfiez-vous des pop-ups « Votre PC est infecté ! » : ce sont souvent EUX, l'arnaque.",
)

slide_cards(
    "Les sauvegardes : votre filet de sécurité",
    [
        ("3 copies", "L'original + 2 sauvegardes de vos données", TEAL),
        ("2 supports", "Ex. disque externe + cloud", TEAL),
        ("1 hors du domicile", "À l'abri du vol, de l'incendie, du ransomware", TEAL),
    ],
    kicker="La règle 3–2–1",
    accent=TEAL,
    cols=3,
)
slide_bullets(
    "Pourquoi la sauvegarde change tout",
    [
        "Rançongiciel, vol, casse, incendie, mauvaise manipulation…",
        "Avec une bonne sauvegarde : vous restaurez et vous continuez",
        "Sans sauvegarde : photos et documents perdus définitivement",
        "Testez de temps en temps que la sauvegarde fonctionne vraiment",
    ],
    kicker="Le réflexe qui sauve",
    accent=TEAL,
    note="Une sauvegarde jamais testée est une sauvegarde en laquelle on ne peut pas avoir confiance.",
)

slide_bullets(
    "Naviguer sur internet en sécurité",
    [
        "Vérifiez l'adresse du site (l'URL) avant de saisir quoi que ce soit",
        "« https » = connexion chiffrée, mais ça ne veut pas dire « site de confiance »",
        "Méfiez-vous des pop-ups, boutons « Télécharger » clignotants, fausses alertes",
        "Sur un ordinateur partagé : déconnectez-vous et fermez la session",
    ],
    kicker="Au quotidien",
    accent=BLUE,
)

slide_bullets(
    "Protéger sa messagerie : la clé du royaume",
    [
        "Votre boîte e-mail sert à réinitialiser TOUS vos autres comptes",
        "La pirater = pouvoir prendre la main sur le reste (banque, réseaux…)",
        "Donc : mot de passe unique et solide + double authentification",
        "Méfiance sur les pièces jointes et les liens, même venant d'un « proche »",
    ],
    kicker="Priorité absolue",
    accent=RED,
    note="Si vous ne sécurisez qu'un seul compte aujourd'hui, que ce soit votre e-mail.",
)

slide_bullets(
    "Réseaux sociaux et vie privée",
    [
        "Réglez la confidentialité : qui voit vos publications ?",
        "Limitez les infos publiques (adresse, lieu de travail, absences, enfants)",
        "Méfiez-vous des petits quiz / jeux qui aspirent vos données",
        "Ce qui est publié en ligne peut y rester… très longtemps",
        "Une photo de vacances en direct annonce aussi : « la maison est vide »",
    ],
    kicker="Maîtriser son image",
    accent=BLUE,
)

# --- PARTIE 4 : CAS CONCRETS --------------------------------------------
slide_section("04", "Cas concrets de la vie quotidienne", "Smartphone, objets connectés, banque, enfants, télétravail")

slide_bullets(
    "Votre smartphone : le plus personnel des appareils",
    [
        "Verrouillez-le : code à 6 chiffres (pas 0000) + empreinte / visage",
        "N'installez des applis QUE depuis les magasins officiels (App Store, Play Store)",
        "Vérifiez les autorisations : pourquoi cette lampe torche veut vos contacts ?",
        "Activez la localisation à distance (« Localiser mon appareil ») en cas de perte",
    ],
    kicker="Cas n°1",
    accent=TEAL,
    note="Un téléphone contient souvent plus d'informations sensibles qu'un ordinateur.",
)

slide_bullets(
    "Les objets connectés (maison intelligente)",
    [
        "Caméras, montres, assistants vocaux, TV, thermostats, jouets…",
        "Changez IMMÉDIATEMENT le mot de passe par défaut (souvent « admin/admin »)",
        "Mettez-les à jour, comme le reste",
        "Chaque objet connecté est une porte d'entrée potentielle dans votre foyer",
    ],
    kicker="Cas n°2",
    accent=ORANGE,
    note="Une caméra de surveillance mal configurée peut… servir à vous surveiller.",
)

slide_bullets(
    "Achats et banque en ligne",
    [
        "Achetez sur des sites connus ; méfiez-vous des prix « trop beaux »",
        "Utilisez une carte virtuelle ou fixez des plafonds si votre banque le permet",
        "Surveillez régulièrement vos relevés bancaires",
        "Votre banque ne vous demandera JAMAIS votre code ou vos codes par téléphone/SMS",
    ],
    kicker="Cas n°3",
    accent=TEAL,
    note="Un doute sur un appel « de votre banque » ? Raccrochez et rappelez le numéro officiel.",
)

slide_bullets(
    "Les enfants et le numérique",
    [
        "Activez le contrôle parental (temps d'écran, contenus, achats)",
        "Privilégiez le dialogue et la confiance à la surveillance pure",
        "Expliquez le harcèlement, les inconnus, les images qu'on ne partage pas",
        "Protégez aussi LEURS données : limitez ce que vous publiez sur eux",
    ],
    kicker="Cas n°4",
    accent=BLUE,
)

slide_bullets(
    "Le télétravail et les usages pro",
    [
        "Séparez autant que possible les usages professionnels et personnels",
        "Utilisez le réseau et le VPN fournis par l'employeur pour les accès sensibles",
        "Ne stockez pas de documents pro sensibles n'importe où",
        "Verrouillez votre session dès que vous quittez votre poste",
    ],
    kicker="Cas n°5",
    accent=TEAL,
)

# --- PARTIE 5 : RÉAGIR ---------------------------------------------------
slide_section("05", "Réagir en cas de problème", "Garder son calme, agir vite, savoir à qui s'adresser")

slide_bullets(
    "Reconnaître qu'il se passe quelque chose",
    [
        "Appareil anormalement lent, surchauffe, pop-ups incessants",
        "Fichiers devenus inaccessibles ou renommés (signe de rançongiciel)",
        "Mots de passe qui ne fonctionnent plus, connexions inconnues signalées",
        "Vos proches reçoivent des messages étranges « de votre part »",
        "Débits bancaires que vous ne reconnaissez pas",
    ],
    kicker="Les signaux",
    accent=ORANGE,
)

slide_checklist(
    "Que faire immédiatement ?",
    [
        "Garder son calme : la panique fait faire des erreurs",
        "Déconnecter l'appareil d'internet (Wi-Fi / câble)",
        "Changer ses mots de passe depuis un AUTRE appareil sain",
        "Prévenir sa banque en cas de fraude, faire opposition",
        "Conserver les preuves (captures d'écran, messages)",
        "Ne jamais payer une rançon sans avis d'un professionnel",
    ],
    kicker="Les bons gestes",
    accent=TEAL,
    cols=2,
    note="Agir vite limite les dégâts — mais agir calmement les limite encore plus.",
)

slide_cards(
    "À qui s'adresser (en France)",
    [
        ("cybermalveillance.gouv.fr", "Aide, diagnostic, mise en relation", TEAL),
        ("17Cyber", "Le guichet unique d'assistance", TEAL),
        ("Déposer plainte", "Commissariat / gendarmerie", BLUE),
        ("Signalement", "internet-signalement.gouv.fr (Pharos)", BLUE),
        ("La CNIL", "Pour vos données personnelles — cnil.fr", ORANGE),
        ("Votre banque", "Opposition et remboursement de fraude", RED),
    ],
    kicker="Vous n'êtes pas seul",
    accent=TEAL,
    cols=3,
)

slide_bullets(
    "Vos droits : le RGPD en une diapo",
    [
        "Le RGPD encadre l'usage de vos données personnelles en Europe",
        ("Il vous donne des droits concrets :", 0),
        ("accès, rectification, effacement (« droit à l'oubli »), opposition", 1),
        "Votre consentement doit être libre et éclairé (cookies, newsletters…)",
        "En cas d'abus : la CNIL est là pour vous",
    ],
    kicker="Le cadre légal",
    accent=ORANGE,
    note="Vous avez le droit de demander à une entreprise quelles données elle détient sur vous.",
)

# --- CONCLUSION ----------------------------------------------------------
slide_section("06", "Conclusion", "L'essentiel à emporter chez soi")

slide_checklist(
    "Les 10 bonnes habitudes",
    [
        "Des mots de passe longs et uniques",
        "Un gestionnaire de mots de passe",
        "La double authentification (surtout l'e-mail)",
        "Les mises à jour activées",
        "Des sauvegardes régulières (règle 3-2-1)",
        "Réfléchir avant de cliquer",
        "Vérifier l'expéditeur et l'adresse (URL)",
        "Prudence sur le Wi-Fi public",
        "Régler la confidentialité de ses comptes",
        "Savoir demander de l'aide sans avoir honte",
    ],
    kicker="Votre check-list",
    accent=TEAL,
    cols=2,
    note="Aucune de ces habitudes n'est technique. Toutes sont à votre portée dès ce soir.",
)

slide_stat(
    "3 secondes",
    "Le réflexe à retenir : dans le doute, on ne clique pas — on vérifie",
    "Ces quelques secondes de recul sont votre meilleure protection, tous les jours.",
    accent=NAVY,
)

slide_cards(
    "Pour aller plus loin (ressources fiables)",
    [
        ("cybermalveillance.gouv.fr", "Conseils et assistance officiels", TEAL),
        ("cnil.fr", "Vie privée et données personnelles", BLUE),
        ("ssi.gouv.fr (ANSSI)", "Guides pratiques de référence", BLUE),
        ("haveibeenpwned.com", "Vérifier si un e-mail a fuité", ORANGE),
        ("Le gestionnaire choisi", "Bitwarden, KeePass…", TEAL),
        ("Cet atelier", "Vos notes + la check-list", GREY),
    ],
    kicker="À garder sous la main",
    accent=BLUE,
    cols=3,
)

slide_bullets(
    "Le message à retenir",
    [
        "La cybersécurité, ce n'est pas de la magie : c'est du bon sens outillé",
        "Quelques réflexes simples éliminent l'immense majorité des risques",
        "On ne cherche pas le risque zéro, mais à ne plus être une cible facile",
        "La sécurité est un travail d'équipe : parlez-en autour de vous",
    ],
    kicker="En un mot",
    accent=TEAL,
    note="Protéger ses données, c'est protéger sa tranquillité, son argent et ses proches.",
)

# Merci / questions
s = new_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(5.5), SW, Inches(0.09), TEAL)
tb, tf = textbox(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.4),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Merci de votre attention !", 46, WHITE, bold=True, first=True,
     align=PP_ALIGN.CENTER, font=FONT_H)
para(tf, "Place à vos questions et à vos expériences", 22,
     RGBColor(0xC7, 0xD6, 0xE6), align=PP_ALIGN.CENTER, space_before=10)

# --------------------------------------------------------------------------
OUT = "cybersecurite_presentation.pptx"
prs.save(OUT)
print(f"OK — {len(prs.slides._sldIdLst)} diapositives écrites dans {OUT}")
