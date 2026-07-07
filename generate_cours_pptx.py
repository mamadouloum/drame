#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère la présentation PowerPoint :
    « Normes et méthodes de sécurité » — cours de licence 3 (40 h)

Design « éditorial premium » : fond papier chaud, accents or / bleu nuit / teal,
diapositives de module sur fond bleu nuit profond.

Le contenu est lu depuis `cours_contenu.SLIDES` (source unique partagée avec le
générateur de discours), ce qui garantit l'alignement diapo par diapo.

Sortie : normes_methodes_securite_40h.pptx (format 16:9)
Dépendances : python-pptx
"""

import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import cours_contenu

# --------------------------------------------------------------------------
# Palette « premium »
# --------------------------------------------------------------------------
INK    = RGBColor(0x0E, 0x1B, 0x2A)   # bleu nuit profond (fonds sombres)
INK2   = RGBColor(0x15, 0x27, 0x3B)   # bleu nuit secondaire
INK3   = RGBColor(0x22, 0x38, 0x50)   # surbrillance sombre
GOLD   = RGBColor(0xC7, 0xA2, 0x4A)   # or élégant (accent principal)
GOLDL  = RGBColor(0xE4, 0xCB, 0x86)   # or clair
TEAL   = RGBColor(0x2F, 0x7E, 0x78)   # teal profond (accent secondaire)
BLUE   = RGBColor(0x35, 0x6A, 0x9A)   # bleu ardoise
RED    = RGBColor(0xB0, 0x3A, 0x2E)   # rouge sobre (danger)
ORANGE = RGBColor(0xC2, 0x6B, 0x2D)   # ambre (vigilance)
PAPER  = RGBColor(0xF6, 0xF4, 0xEF)   # papier chaud (fond de contenu)
PAPER2 = RGBColor(0xEF, 0xEB, 0xE1)   # papier plus soutenu
CARD   = RGBColor(0xFF, 0xFF, 0xFF)   # cartes
LINE   = RGBColor(0xDD, 0xD6, 0xC7)   # filets / bordures
INKTX  = RGBColor(0x1C, 0x27, 0x33)   # texte principal
GREY   = RGBColor(0x63, 0x6F, 0x7B)   # texte secondaire
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xCE, 0xD8, 0xE2)   # texte clair sur fond sombre

TOKENS = {"GOLD": GOLD, "GOLDL": GOLDL, "TEAL": TEAL, "BLUE": BLUE, "RED": RED,
          "ORANGE": ORANGE, "GREY": GREY, "INK": INK, "INK2": INK2, "WHITE": WHITE}

FONT   = "Calibri"
FONT_H = "Georgia"          # empattement pour les titres = touche « éditoriale »

COURSE = "Normes et méthodes de sécurité"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def tok(name, default=GOLD):
    if name is None:
        return default
    if isinstance(name, RGBColor):
        return name
    return TOKENS.get(name, default)


# --------------------------------------------------------------------------
# Primitives
# --------------------------------------------------------------------------
def new_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, l, t, w, h, color, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def rounded(slide, l, t, w, h, color, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def line_h(slide, l, t, w, color, weight=Pt(2.5)):
    rect(slide, l, t, w, weight, color)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         first=False, space_after=6, space_before=0, font=FONT, line_spacing=1.05,
         spacing_pts=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def footer(slide, module=0, num=0):
    tag = "Introduction" if module == 0 else "Module %d" % module
    tb, tf = textbox(slide, Inches(0.55), Inches(7.04), Inches(9), Inches(0.4))
    para(tf, "%s  ·  %s" % (COURSE, tag), 8.5, GREY, first=True)
    tb2, tf2 = textbox(slide, Inches(11.4), Inches(7.0), Inches(1.4), Inches(0.4))
    para(tf2, "%02d" % num, 10, GOLD, align=PP_ALIGN.RIGHT, first=True, bold=True)


def page_number(slide, num):
    """Numéro discret pour les diapos sur fond sombre."""
    tb, tf = textbox(slide, Inches(11.4), Inches(7.0), Inches(1.4), Inches(0.4))
    para(tf, "%02d" % num, 10, GOLD, align=PP_ALIGN.RIGHT, first=True, bold=True)


def content_header(slide, title, kicker=None, accent=GOLD):
    rect(slide, 0, 0, SW, SH, PAPER)
    rect(slide, 0, 0, SW, Inches(0.16), accent)        # fine barre couleur
    top = Inches(0.52)
    if kicker:
        tb, tf = textbox(slide, Inches(0.62), Inches(0.44), Inches(11.5), Inches(0.4))
        para(tf, kicker.upper(), 12, accent, bold=True, first=True, font=FONT)
        top = Inches(0.86)
    tb, tf = textbox(slide, Inches(0.6), top, Inches(12.2), Inches(0.95))
    para(tf, title, 29, INK, bold=True, first=True, font=FONT_H, line_spacing=1.0)
    line_h(slide, Inches(0.64), top + Inches(0.9), Inches(1.5), accent, Pt(3))


# --------------------------------------------------------------------------
# Puces (partagées)
# --------------------------------------------------------------------------
def bullets(tf, items, size=18, color=INKTX, accent=GOLD, first_para=True,
            space_after=10):
    started = not first_para
    for it in items:
        text, lvl = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if (first_para and not started) else tf.add_paragraph()
        started = True
        p.space_after = Pt(space_after if lvl == 0 else 4)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        indent = "" if lvl == 0 else "      "
        mk = p.add_run()
        mk.text = indent + ("◆  " if lvl == 0 else "–  ")
        mk.font.size = Pt(size - 3 if lvl == 0 else size - 4)
        mk.font.color.rgb = accent if lvl == 0 else GREY
        mk.font.bold = True
        mk.font.name = FONT
        rn = p.add_run()
        rn.text = text
        rn.font.size = Pt(size if lvl == 0 else size - 3)
        rn.font.color.rgb = color if lvl == 0 else GREY
        rn.font.name = FONT


def note_band(slide, text, accent=GOLD, y=6.5):
    rounded(slide, Inches(0.6), Inches(y), Inches(12.13), Inches(0.62), PAPER2)
    rect(slide, Inches(0.6), Inches(y), Inches(0.1), Inches(0.62), accent)
    tb, tf = textbox(slide, Inches(0.9), Inches(y), Inches(11.6), Inches(0.62),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, 13, INK, italic=True, first=True)


# ==========================================================================
#  GABARITS
# ==========================================================================
def r_title(s):
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, 0, Inches(0.28), SH, GOLD)
    # bandes décoratives fines
    rect(s, Inches(0.9), Inches(2.02), Inches(2.2), Pt(3), GOLD)
    tb, tf = textbox(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.5))
    para(tf, "COURS DE LICENCE 3", 13, GOLDL, bold=True, first=True, font=FONT)
    tb, tf = textbox(s, Inches(0.86), Inches(2.25), Inches(11.6), Inches(2.9))
    for i, part in enumerate(s._title.split("\n")):
        para(tf, part, 50, WHITE, bold=True, first=(i == 0), font=FONT_H,
             line_spacing=1.0, space_after=0)
    tb, tf = textbox(s, Inches(0.9), Inches(4.75), Inches(11.4), Inches(1.0))
    para(tf, s._d["subtitle"], 21, CREAM, first=True, font=FONT, line_spacing=1.1)
    rect(s, 0, Inches(6.35), SW, Pt(1.2), INK3)
    tb, tf = textbox(s, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.7))
    para(tf, s._d["meta"], 13, GOLDL, first=True, font=FONT)
    page_number(s, s._num)


def r_module(s):
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, 0, Inches(0.28), SH, GOLD)
    # grand numéro
    tb, tf = textbox(s, Inches(0.8), Inches(0.75), Inches(6), Inches(2.2))
    para(tf, s._d["number"], 130, INK3, bold=True, first=True, font=FONT_H,
         line_spacing=0.9)
    tb, tf = textbox(s, Inches(0.86), Inches(0.9), Inches(6), Inches(1.0))
    para(tf, "MODULE", 16, GOLD, bold=True, first=True, font=FONT)
    # titre
    tb, tf = textbox(s, Inches(0.9), Inches(3.05), Inches(11.6), Inches(1.7))
    para(tf, s._title, 38, WHITE, bold=True, first=True, font=FONT_H, line_spacing=1.0)
    line_h(s, Inches(0.94), Inches(3.0), Inches(1.8), GOLD, Pt(3))
    if s._d.get("subtitle"):
        tb, tf = textbox(s, Inches(0.9), Inches(4.55), Inches(11.4), Inches(0.6))
        para(tf, s._d["subtitle"], 19, GOLDL, italic=True, first=True)
    # sujets abordés
    topics = s._d.get("topics") or []
    if topics:
        y = 5.25
        tb, tf = textbox(s, Inches(0.92), Inches(y), Inches(11.5), Inches(1.9))
        para(tf, "AU PROGRAMME", 11, TEAL, bold=True, first=True, space_after=6)
        for tpc in topics:
            p = tf.add_paragraph()
            p.space_after = Pt(3)
            mk = p.add_run(); mk.text = "◆  "
            mk.font.size = Pt(12); mk.font.color.rgb = GOLD; mk.font.name = FONT
            rn = p.add_run(); rn.text = tpc
            rn.font.size = Pt(14); rn.font.color.rgb = CREAM; rn.font.name = FONT
    page_number(s, s._num)


def r_objectives(s):
    accent = tok(s._d.get("accent"), GOLD)
    content_header(s, s._title, s._d.get("kicker"), accent)
    items = s._d["items"]
    top0 = 2.0
    line = 0.66
    for i, it in enumerate(items):
        t = Inches(top0 + i * line)
        badge = rounded(s, Inches(0.7), t, Inches(0.44), Inches(0.44), accent)
        tbb, tfb = textbox(s, Inches(0.7), t - Inches(0.02), Inches(0.44), Inches(0.46),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tfb, "✓", 16, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
        tb, tf = textbox(s, Inches(1.32), t - Inches(0.04), Inches(11.3), Inches(0.55),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, it, 17, INKTX, first=True, line_spacing=1.0)
    footer(s, s._module, s._num)


def r_bullets(s):
    accent = tok(s._d.get("accent"), GOLD)
    content_header(s, s._title, s._d.get("kicker"), accent)
    note = s._d.get("note")
    h = Inches(3.9) if note else Inches(4.7)
    tb, tf = textbox(s, Inches(0.72), Inches(1.95), Inches(11.9), h)
    bullets(tf, s._d["items"], size=18, accent=accent)
    if note:
        note_band(s, note, accent, y=6.35)
    footer(s, s._module, s._num)


def r_cards(s):
    accent = tok(s._d.get("accent"), GOLD)
    content_header(s, s._title, s._d.get("kicker"), accent)
    cards = s._d["cards"]
    cols = s._d.get("cols", 3)
    note = s._d.get("note")
    margin = Inches(0.6)
    gutter = Inches(0.28)
    total_w = SW - 2 * margin
    card_w = int((total_w - (cols - 1) * gutter) / cols)
    rows = math.ceil(len(cards) / cols)
    top0 = Inches(2.0)
    avail_h = Inches(3.95) if note else Inches(4.75)
    card_h = int((avail_h - (rows - 1) * gutter) / rows)
    for i, c in enumerate(cards):
        if len(c) == 3:
            ctitle, desc, col = c[0], c[1], tok(c[2], accent)
        else:
            ctitle, desc = c[0], c[1]; col = accent
        r = i // cols
        cc = i % cols
        l = margin + cc * (card_w + gutter)
        t = top0 + r * (card_h + gutter)
        rounded(s, l, t, card_w, card_h, CARD, LINE)
        rect(s, l, t, Inches(0.11), card_h, col)
        tb, tf = textbox(s, l + Inches(0.3), t + Inches(0.16),
                         card_w - Inches(0.45), card_h - Inches(0.28))
        para(tf, ctitle, 15.5, INK, bold=True, first=True, space_after=4, font=FONT_H)
        if desc:
            para(tf, desc, 12.5, GREY, line_spacing=1.02)
    if note:
        note_band(s, note, accent, y=6.35)
    footer(s, s._module, s._num)


def r_twocol(s):
    accent = tok(s._d.get("accent"), GOLD)
    content_header(s, s._title, s._d.get("kicker"), accent)
    lc = tok(s._d.get("left_color"), RED)
    rc = tok(s._d.get("right_color"), TEAL)
    top = Inches(2.0)
    colw = Inches(5.9)
    gap = Inches(0.33)
    ch = Inches(4.7)
    # gauche
    rounded(s, Inches(0.6), top, colw, ch, CARD, LINE)
    rect(s, Inches(0.6), top, colw, Inches(0.6), lc)
    tb, tf = textbox(s, Inches(0.82), top + Inches(0.04), colw - Inches(0.4),
                     Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, s._d["left_head"], 17, WHITE, bold=True, first=True)
    tb, tf = textbox(s, Inches(0.85), top + Inches(0.78), colw - Inches(0.5), Inches(3.7))
    bullets(tf, s._d["left_items"], size=15, accent=lc, space_after=8)
    # droite
    lx = Inches(0.6) + colw + gap
    rounded(s, lx, top, colw, ch, CARD, LINE)
    rect(s, lx, top, colw, Inches(0.6), rc)
    tb, tf = textbox(s, lx + Inches(0.22), top + Inches(0.04), colw - Inches(0.4),
                     Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, s._d["right_head"], 17, WHITE, bold=True, first=True)
    tb, tf = textbox(s, lx + Inches(0.25), top + Inches(0.78), colw - Inches(0.5), Inches(3.7))
    bullets(tf, s._d["right_items"], size=15, accent=rc, space_after=8)
    footer(s, s._module, s._num)


def r_matrix(s):
    accent = tok(s._d.get("accent"), GOLD)
    content_header(s, s._title, s._d.get("kicker"), accent)
    headers = s._d["headers"]
    rows = s._d["rows"]
    note = s._d.get("note")
    ncol = len(headers)
    left = Inches(0.6)
    total_w = SW - Inches(1.2)
    # largeurs : 1re colonne plus étroite si 2 colonnes
    if ncol == 2:
        widths = [int(total_w * 0.32), int(total_w * 0.68)]
    else:
        widths = [int(total_w / ncol)] * ncol
    top = Inches(2.0)
    hrow = Inches(0.56)
    body_h = (Inches(3.7) if note else Inches(4.5))
    rh = int(body_h / max(len(rows), 1))
    rh = min(rh, Inches(0.7))
    # en-tête
    x = left
    for j, htxt in enumerate(headers):
        rect(s, x, top, widths[j], hrow, INK)
        tb, tf = textbox(s, x + Inches(0.16), top, widths[j] - Inches(0.3), hrow,
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, htxt, 13.5, WHITE, bold=True, first=True)
        x += widths[j]
    # lignes
    for i, row in enumerate(rows):
        y = top + hrow + i * rh
        bg = CARD if i % 2 == 0 else PAPER2
        x = left
        for j, cell in enumerate(row):
            rect(s, x, y, widths[j], rh, bg, LINE, Pt(0.5))
            tb, tf = textbox(s, x + Inches(0.16), y, widths[j] - Inches(0.3), rh,
                             anchor=MSO_ANCHOR.MIDDLE)
            para(tf, cell, 12.5, INKTX if j else INK,
                 bold=(j == 0), first=True, line_spacing=1.0)
            x += widths[j]
    if note:
        note_band(s, note, accent, y=6.4)
    footer(s, s._module, s._num)


def r_process(s):
    accent = tok(s._d.get("accent"), GOLD)
    content_header(s, s._title, s._d.get("kicker"), accent)
    steps = s._d["steps"]
    n = len(steps)
    top0 = 2.05
    gap = 0.2
    avail = 4.6
    bh = (avail - (n - 1) * gap) / n
    for i, (t_, d_) in enumerate(steps):
        y = Inches(top0 + i * (bh + gap))
        rounded(s, Inches(0.6), y, Inches(12.13), Inches(bh), CARD, LINE)
        # badge numéro
        rect(s, Inches(0.6), y, Inches(1.25), Inches(bh), INK)
        tbb, tfb = textbox(s, Inches(0.6), y, Inches(1.25), Inches(bh),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tfb, str(i + 1), 30, GOLD, bold=True, first=True, align=PP_ALIGN.CENTER,
             font=FONT_H)
        tb, tf = textbox(s, Inches(2.1), y, Inches(10.4), Inches(bh),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t_, 17, INK, bold=True, first=True, space_after=2, font=FONT_H)
        para(tf, d_, 13.5, GREY, line_spacing=1.0)
        # flèche entre étapes
        if i < n - 1:
            tbb, tfb = textbox(s, Inches(1.0), y + Inches(bh) - Inches(0.08),
                               Inches(0.5), Inches(0.3))
            para(tfb, "▼", 11, accent, first=True, align=PP_ALIGN.CENTER)
    footer(s, s._module, s._num)


def r_stat(s):
    accent = tok(s._d.get("accent"), TEAL)
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, 0, SW, Inches(0.16), accent)
    rect(s, 0, SH - Inches(0.16), SW, Inches(0.16), accent)
    big = s._d.get("big", s._title)
    tb, tf = textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.1),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, 100, GOLD, bold=True, first=True, align=PP_ALIGN.CENTER,
         font=FONT_H, line_spacing=0.9)
    tb, tf = textbox(s, Inches(1.1), Inches(4.35), Inches(11.1), Inches(1.5))
    para(tf, s._d["caption"], 24, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         line_spacing=1.08, font=FONT_H)
    if s._d.get("sub"):
        tb, tf = textbox(s, Inches(1.6), Inches(5.75), Inches(10.1), Inches(1.0))
        para(tf, s._d["sub"], 14, CREAM, first=True, align=PP_ALIGN.CENTER,
             italic=True, line_spacing=1.05)
    page_number(s, s._num)


def r_key(s):
    accent = tok(s._d.get("accent"), GOLD)
    rect(s, 0, 0, SW, SH, INK2)
    rect(s, 0, 0, Inches(0.28), SH, accent)
    tb, tf = textbox(s, Inches(1.2), Inches(0.8), Inches(6), Inches(0.5))
    para(tf, (s._d.get("kicker") or "À RETENIR").upper(), 13, accent, bold=True,
         first=True, font=FONT)
    # grand guillemet
    tb, tf = textbox(s, Inches(1.0), Inches(1.1), Inches(2), Inches(1.4))
    para(tf, "“", 90, INK3, bold=True, first=True, font=FONT_H)
    tb, tf = textbox(s, Inches(1.3), Inches(2.35), Inches(10.8), Inches(3.4),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, s._d["text"], 27, WHITE, first=True, font=FONT_H, line_spacing=1.15)
    if s._d.get("attrib"):
        line_h(s, Inches(1.34), Inches(5.95), Inches(1.2), accent, Pt(2.5))
        tb, tf = textbox(s, Inches(1.32), Inches(6.05), Inches(10), Inches(0.5))
        para(tf, s._d["attrib"], 14, GOLDL, first=True, italic=True)
    page_number(s, s._num)


def r_exercise(s):
    accent = GOLD
    rect(s, 0, 0, SW, SH, PAPER)
    rect(s, 0, 0, SW, Inches(0.16), accent)
    # bandeau TD
    rounded(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.5), INK)
    tb, tf = textbox(s, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.5),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "✎  " + (s._d.get("kicker") or "EXERCICE"), 12, GOLD, bold=True,
         first=True, align=PP_ALIGN.CENTER)
    if s._d.get("duration"):
        rounded(s, Inches(11.1), Inches(0.5), Inches(1.63), Inches(0.5), TEAL)
        tb, tf = textbox(s, Inches(11.1), Inches(0.5), Inches(1.63), Inches(0.5),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "⏱  " + s._d["duration"], 12, WHITE, bold=True, first=True,
             align=PP_ALIGN.CENTER)
    tb, tf = textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.9))
    para(tf, s._title, 28, INK, bold=True, first=True, font=FONT_H)
    tb, tf = textbox(s, Inches(0.62), Inches(2.0), Inches(12.0), Inches(0.9))
    para(tf, s._d["brief"], 15, INKTX, first=True, italic=True, line_spacing=1.1)
    # tâches
    tasks = s._d["tasks"]
    top0 = 3.0
    line = 0.62
    for i, tk in enumerate(tasks):
        t = Inches(top0 + i * line)
        badge = rounded(s, Inches(0.7), t, Inches(0.42), Inches(0.42), accent)
        tbb, tfb = textbox(s, Inches(0.7), t - Inches(0.02), Inches(0.42), Inches(0.46),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tfb, str(i + 1), 14, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
        tb, tf = textbox(s, Inches(1.3), t - Inches(0.04), Inches(11.3), Inches(0.55),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, tk, 15.5, INKTX, first=True, line_spacing=1.0)
    if s._d.get("deliverable"):
        y = 3.0 + len(tasks) * line + 0.1
        rounded(s, Inches(0.6), Inches(y), Inches(12.13), Inches(0.62), INK)
        tb, tf = textbox(s, Inches(0.85), Inches(y), Inches(11.7), Inches(0.62),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "Livrable attendu : " + s._d["deliverable"], 13.5, GOLDL,
             bold=True, first=True)
    footer(s, s._module, s._num)


def r_checklist(s):
    accent = tok(s._d.get("accent"), TEAL)
    content_header(s, s._title, s._d.get("kicker"), accent)
    items = s._d["items"]
    cols = s._d.get("cols", 2)
    note = s._d.get("note")
    margin = 0.7
    colw = 5.95
    gap = 0.35
    per_col = math.ceil(len(items) / cols)
    top0 = 2.05
    line = 0.66
    for i, it in enumerate(items):
        col = i // per_col
        row = i % per_col
        l = Inches(margin + col * (colw + gap))
        t = Inches(top0 + row * line)
        rounded(s, l, t, Inches(0.42), Inches(0.42), accent)
        tbb, tfb = textbox(s, l, t - Inches(0.02), Inches(0.42), Inches(0.46),
                           anchor=MSO_ANCHOR.MIDDLE)
        para(tfb, "✓", 14, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
        tb, tf = textbox(s, l + Inches(0.55), t - Inches(0.05), Inches(colw - 0.55),
                         Inches(0.56), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, it, 14.5, INKTX, first=True, line_spacing=1.0)
    if note:
        note_band(s, note, accent, y=6.4)
    footer(s, s._module, s._num)


def r_quiz(s):
    accent = tok(s._d.get("accent"), ORANGE)
    content_header(s, s._title, s._d.get("kicker", "ON VÉRIFIE"), accent)
    qa = s._d["qa"]
    top = 2.0
    n = len(qa)
    gap = 0.2
    bh = (4.7 - (n - 1) * gap) / n
    for i, (q, a) in enumerate(qa):
        t = Inches(top + i * (bh + gap))
        rounded(s, Inches(0.62), t, Inches(12.1), Inches(bh), CARD, LINE)
        rect(s, Inches(0.62), t, Inches(0.11), Inches(bh), accent)
        tb, tf = textbox(s, Inches(0.95), t + Inches(0.1), Inches(11.5), Inches(bh) - Inches(0.15),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "?  " + q, 15, INK, bold=True, first=True, space_after=3)
        para(tf, "→  " + a, 13, TEAL, italic=True, line_spacing=1.0)
    footer(s, s._module, s._num)


def r_closing(s):
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, Inches(5.5), SW, Pt(3), GOLD)
    tb, tf = textbox(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.4),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, s._title, 44, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         font=FONT_H)
    if s._d.get("subtitle"):
        para(tf, s._d["subtitle"], 20, CREAM, align=PP_ALIGN.CENTER, space_before=12)
    page_number(s, s._num)


RENDERERS = {
    "title": r_title, "module": r_module, "objectives": r_objectives,
    "bullets": r_bullets, "cards": r_cards, "twocol": r_twocol,
    "matrix": r_matrix, "process": r_process, "stat": r_stat, "key": r_key,
    "exercise": r_exercise, "checklist": r_checklist, "quiz": r_quiz,
    "closing": r_closing,
}


class _SlideProxy:
    """Expose l'API slide (shapes) et les champs de contenu (_d, _title, _module)
    attendus par les fonctions de gabarit."""
    def __init__(self, slide, d, num):
        self._slide = slide
        self._d = d
        self._title = d.get("title", "")
        self._module = d.get("module", 0)
        self._num = num

    @property
    def shapes(self):
        return self._slide.shapes


def build():
    for i, d in enumerate(cours_contenu.SLIDES):
        slide = new_slide()
        RENDERERS[d["kind"]](_SlideProxy(slide, d, i + 1))
    out = "normes_methodes_securite_40h.pptx"
    prs.save(out)
    print("OK — %d diapositives écrites dans %s" % (len(prs.slides._sldIdLst), out))


if __name__ == "__main__":
    build()
